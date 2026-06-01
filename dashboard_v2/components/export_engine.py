"""
Export Engine · Generate PPTX, PDF, and Excel reports from KPI data.
"""

import io
import datetime
import pandas as pd
import numpy as np
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from fpdf import FPDF
import xlsxwriter
from data.loader import (
    apply_filters, get_parieurs, get_utilisateurs_betclic,
    calc_tom, calc_notoriete_totale, calc_penetration, calc_consideration,
    calc_preference, calc_wallet_share, calc_rappel_campagne,
    calc_satisfaction, calc_nps,
    calc_image_scores, calc_notoriete_all_brands, calc_tom_all_brands,
    calc_marque_principale_all, calc_funnel, calc_kpi_by_vague,
    calc_delta, get_latest_vague, get_previous_vague,
    COMPETITORS, IMAGE_ATTRIBUTES, VAGUE_SHORT, CITIES
)

BETCLIC_RED = RGBColor(0xC0, 0x39, 0x2B)
OW_PURPLE = RGBColor(0x6C, 0x34, 0x83)
SLATE_RGB = RGBColor(0x2C, 0x3E, 0x50)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xF5, 0xF6, 0xF8)


# ═══════════════════════════════════════════
# COMMON KPI COMPUTATION
# ═══════════════════════════════════════════

def _compute_all_kpis(df, vagues_list):
    """Compute all key KPIs for the filtered dataset, per vague and latest."""
    kpis = {}

    # Per-vague
    kpis["tom_vagues"] = calc_kpi_by_vague(df, calc_tom)
    kpis["notoriete_vagues"] = calc_kpi_by_vague(df, calc_notoriete_totale)
    kpis["penetration_vagues"] = calc_kpi_by_vague(df, calc_penetration)
    kpis["consideration_vagues"] = calc_kpi_by_vague(df, calc_consideration)
    kpis["preference_vagues"] = calc_kpi_by_vague(df, calc_preference)
    kpis["wallet_vagues"] = calc_kpi_by_vague(df, calc_wallet_share)
    kpis["rappel_vagues"] = calc_kpi_by_vague(df, calc_rappel_campagne)
    kpis["satisfaction_vagues"] = calc_kpi_by_vague(df, calc_satisfaction)

    def _nps_score(d):
        return calc_nps(d)["nps"]
    kpis["nps_vagues"] = calc_kpi_by_vague(df, _nps_score)

    # Latest vague (dynamic: handles V1 only, V1+V2, etc.)
    _candidates = [v for v in vagues_list if v in df["Vague"].unique()]
    if _candidates:
        latest_v = sorted(
            _candidates,
            key=lambda x: int(x.replace("Vague ", "")) if x.replace("Vague ", "").isdigit() else 999,
        )[-1]
    else:
        latest_v = "Vague 1"
    df_latest = df[df["Vague"] == latest_v]
    kpis["latest_vague"] = latest_v

    kpis["tom"] = calc_tom(df_latest)
    kpis["notoriete"] = calc_notoriete_totale(df_latest)
    kpis["penetration"] = calc_penetration(df_latest)
    kpis["consideration"] = calc_consideration(df_latest)
    kpis["preference"] = calc_preference(df_latest)
    kpis["wallet"] = calc_wallet_share(df_latest)
    kpis["rappel"] = calc_rappel_campagne(df_latest)
    kpis["satisfaction"] = calc_satisfaction(df_latest)
    kpis["nps"] = calc_nps(df_latest)
    kpis["image"] = calc_image_scores(df_latest)
    kpis["notoriete_brands"] = calc_notoriete_all_brands(df_latest)
    kpis["tom_brands"] = calc_tom_all_brands(df_latest)
    kpis["funnel"] = calc_funnel(df_latest)

    # Deltas
    for key in ["tom", "notoriete", "penetration", "consideration", "preference",
                 "wallet", "rappel", "satisfaction"]:
        vague_key = f"{key}_vagues"
        latest_val = get_latest_vague(kpis[vague_key])
        prev_val = get_previous_vague(kpis[vague_key])
        d, direction = calc_delta(latest_val, prev_val)
        kpis[f"{key}_delta"] = d
        kpis[f"{key}_dir"] = direction

    nps_latest = get_latest_vague(kpis["nps_vagues"])
    nps_prev = get_previous_vague(kpis["nps_vagues"])
    d, direction = calc_delta(nps_latest, nps_prev)
    kpis["nps_delta"] = d
    kpis["nps_dir"] = direction

    kpis["n_total"] = len(df_latest)

    return kpis


# ═══════════════════════════════════════════
# PPTX EXPORT
# ═══════════════════════════════════════════

def _add_title_slide(prs, kpis):
    """Add a title slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    # Red top bar
    shape = slide.shapes.add_shape(1, 0, 0, prs.slide_width, Emu(600000))
    shape.fill.solid()
    shape.fill.fore_color.rgb = BETCLIC_RED
    shape.line.fill.background()

    # Title
    txBox = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(1.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "BETCLIC BRAND PULSE"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = SLATE_RGB
    p.alignment = PP_ALIGN.CENTER

    p2 = tf.add_paragraph()
    p2.text = f"Scorecard Exécutif · {kpis['latest_vague']}"
    p2.font.size = Pt(18)
    p2.font.color.rgb = OW_PURPLE
    p2.alignment = PP_ALIGN.CENTER

    p3 = tf.add_paragraph()
    p3.text = f"Base : {kpis['n_total']} répondants"
    p3.font.size = Pt(12)
    p3.font.color.rgb = RGBColor(0x4A, 0x55, 0x68)
    p3.alignment = PP_ALIGN.CENTER

    # Footer
    txBox2 = slide.shapes.add_textbox(Inches(1), Inches(6.5), Inches(8), Inches(0.4))
    tf2 = txBox2.text_frame
    p_f = tf2.paragraphs[0]
    p_f.text = f"© {datetime.date.today().year} OpinionWay Africa × Betclic CI · Confidentiel"
    p_f.font.size = Pt(9)
    p_f.font.color.rgb = RGBColor(0x4A, 0x55, 0x68)
    p_f.alignment = PP_ALIGN.CENTER


def _add_kpi_slide(prs, kpis):
    """Add KPI scorecard slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Header
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "Indicateurs Clés de Performance"
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = BETCLIC_RED

    kpi_list = [
        ("Top-of-Mind", f"{kpis['tom']}%", kpis['tom_delta']),
        ("Notoriété Totale", f"{kpis['notoriete']}%", kpis['notoriete_delta']),
        ("Pénétration Parieurs", f"{kpis['penetration']}%", kpis['penetration_delta']),
        ("Considération", f"{kpis['consideration']}%", kpis['consideration_delta']),
        ("Préférence", f"{kpis['preference']}%", kpis['preference_delta']),
        ("Wallet Share", f"{kpis['wallet']}%", kpis['wallet_delta']),
        ("NPS Score", f"{kpis['nps']['nps']} pts", kpis['nps_delta']),
        ("Satisfaction", f"{kpis['satisfaction']}/5", kpis['satisfaction_delta']),
        ("Rappel Campagne", f"{kpis['rappel']}%", kpis['rappel_delta']),
    ]

    # Build table
    rows = len(kpi_list) + 1
    tbl_shape = slide.shapes.add_table(rows, 3, Inches(0.8), Inches(1.2), Inches(8), Emu(rows * 400000))
    tbl = tbl_shape.table

    # Header row
    for i, header in enumerate(["Indicateur", "Valeur", "Δ vs Vague préc."]):
        cell = tbl.cell(0, i)
        cell.text = header
        for paragraph in cell.text_frame.paragraphs:
            paragraph.font.size = Pt(11)
            paragraph.font.bold = True
            paragraph.font.color.rgb = WHITE
        cell.fill.solid()
        cell.fill.fore_color.rgb = BETCLIC_RED

    for r, (label, value, delta) in enumerate(kpi_list, 1):
        tbl.cell(r, 0).text = label
        tbl.cell(r, 1).text = value
        tbl.cell(r, 2).text = str(delta) if delta else "—"

        for c in range(3):
            cell = tbl.cell(r, c)
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(10)
                paragraph.font.color.rgb = SLATE_RGB
            if r % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = LIGHT_GRAY


def _add_brand_slide(prs, kpis):
    """Add brand awareness comparison slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "Notoriété par Marque"
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = OW_PURPLE

    brands_data = kpis["notoriete_brands"]
    sorted_brands = sorted(brands_data.items(), key=lambda x: x[1], reverse=True)

    rows = len(sorted_brands) + 1
    tbl_shape = slide.shapes.add_table(rows, 3, Inches(1), Inches(1.2), Inches(7.5), Emu(rows * 380000))
    tbl = tbl_shape.table

    for i, header in enumerate(["Marque", "Notoriété Totale", "TOM"]):
        cell = tbl.cell(0, i)
        cell.text = header
        for paragraph in cell.text_frame.paragraphs:
            paragraph.font.size = Pt(11)
            paragraph.font.bold = True
            paragraph.font.color.rgb = WHITE
        cell.fill.solid()
        cell.fill.fore_color.rgb = OW_PURPLE

    for r, (brand, not_val) in enumerate(sorted_brands, 1):
        tom_val = kpis["tom_brands"].get(brand, 0)
        tbl.cell(r, 0).text = brand
        tbl.cell(r, 1).text = f"{not_val}%"
        tbl.cell(r, 2).text = f"{tom_val}%"
        for c in range(3):
            cell = tbl.cell(r, c)
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(10)
                paragraph.font.color.rgb = SLATE_RGB
                if brand == "Betclic":
                    paragraph.font.bold = True
                    paragraph.font.color.rgb = BETCLIC_RED


def _add_image_slide(prs, kpis):
    """Add image attributes slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "Attributs Fonctionnels · Scores Moyens (/5)"
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = BETCLIC_RED

    image_data = kpis["image"]
    sorted_attrs = sorted(image_data.items(), key=lambda x: x[1], reverse=True)

    rows = len(sorted_attrs) + 1
    tbl_shape = slide.shapes.add_table(rows, 2, Inches(1.5), Inches(1.2), Inches(6.5), Emu(rows * 360000))
    tbl = tbl_shape.table

    for i, header in enumerate(["Attribut", "Score moyen"]):
        cell = tbl.cell(0, i)
        cell.text = header
        for paragraph in cell.text_frame.paragraphs:
            paragraph.font.size = Pt(11)
            paragraph.font.bold = True
            paragraph.font.color.rgb = WHITE
        cell.fill.solid()
        cell.fill.fore_color.rgb = BETCLIC_RED

    for r, (attr, score) in enumerate(sorted_attrs, 1):
        tbl.cell(r, 0).text = attr
        tbl.cell(r, 1).text = f"{score}/5"
        for c in range(2):
            for paragraph in tbl.cell(r, c).text_frame.paragraphs:
                paragraph.font.size = Pt(10)
                paragraph.font.color.rgb = SLATE_RGB


def generate_pptx(df, vagues_list) -> bytes:
    """Generate a complete PPTX report."""
    kpis = _compute_all_kpis(df, vagues_list)

    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    _add_title_slide(prs, kpis)
    _add_kpi_slide(prs, kpis)
    _add_brand_slide(prs, kpis)
    _add_image_slide(prs, kpis)

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf.getvalue()


# ═══════════════════════════════════════════
# PDF EXPORT
# ═══════════════════════════════════════════

def _safe(text: str) -> str:
    """Replace unicode chars that Helvetica cannot encode."""
    return (str(text)
            .replace("\u2014", "-").replace("\u2013", "-")
            .replace("\u2019", "'").replace("\u2018", "'")
            .replace("\u201c", '"').replace("\u201d", '"')
            .replace("\u00e9", "e").replace("\u00e8", "e").replace("\u00ea", "e")
            .replace("\u00e0", "a").replace("\u00e2", "a")
            .replace("\u00f4", "o").replace("\u00fb", "u")
            .replace("\u00ee", "i").replace("\u00ef", "i")
            .replace("\u00c9", "E").replace("\u00c8", "E")
            .replace("\u00e7", "c").replace("\u00f9", "u")
            .replace("\u00b0", "o"))


class BetclicPDF(FPDF):
    def header(self):
        self.set_fill_color(192, 57, 43)
        self.rect(0, 0, 210, 8, 'F')
        self.set_font("Helvetica", "B", 9)
        self.set_text_color(192, 57, 43)
        self.set_y(12)
        self.cell(0, 5, "BETCLIC BRAND PULSE - Business Intelligence 2026", align="C")
        self.ln(10)

    def footer(self):
        self.set_y(-15)
        self.set_font("Helvetica", "", 7)
        self.set_text_color(74, 85, 104)
        self.cell(0, 10, f"OpinionWay Africa x Betclic CI - Confidentiel - Page {self.page_no()}/{{nb}}", align="C")

    def section_title(self, title):
        self.set_font("Helvetica", "B", 14)
        self.set_text_color(192, 57, 43)
        self.cell(0, 10, _safe(title), new_x="LMARGIN", new_y="NEXT")
        self.set_draw_color(192, 57, 43)
        self.line(10, self.get_y(), 200, self.get_y())
        self.ln(4)

    def kpi_row(self, label, value, delta=""):
        self.set_font("Helvetica", "", 10)
        self.set_text_color(44, 62, 80)
        self.cell(80, 7, _safe(label), border=0)
        self.set_font("Helvetica", "B", 10)
        self.cell(40, 7, _safe(str(value)), border=0)
        self.set_font("Helvetica", "", 9)
        if delta and str(delta).startswith("+"):
            self.set_text_color(29, 131, 72)
        elif delta and str(delta).startswith("-"):
            self.set_text_color(192, 57, 43)
        else:
            self.set_text_color(212, 118, 10)
        self.cell(40, 7, _safe(str(delta)) if delta else "", border=0, new_x="LMARGIN", new_y="NEXT")
        self.set_text_color(44, 62, 80)


def generate_pdf(df, vagues_list) -> bytes:
    """Generate a PDF report."""
    kpis = _compute_all_kpis(df, vagues_list)

    pdf = BetclicPDF()
    pdf.alias_nb_pages()
    pdf.set_auto_page_break(auto=True, margin=20)

    # Page 1: Cover
    pdf.add_page()
    pdf.ln(30)
    pdf.set_font("Helvetica", "B", 28)
    pdf.set_text_color(44, 62, 80)
    pdf.cell(0, 15, "BETCLIC BRAND PULSE", align="C", new_x="LMARGIN", new_y="NEXT")
    pdf.set_font("Helvetica", "", 16)
    pdf.set_text_color(108, 52, 131)
    pdf.cell(0, 10, f"Scorecard Executif - {kpis['latest_vague']}", align="C", new_x="LMARGIN", new_y="NEXT")
    pdf.ln(5)
    pdf.set_font("Helvetica", "", 11)
    pdf.set_text_color(74, 85, 104)
    pdf.cell(0, 8, f"Base : {kpis['n_total']} repondants", align="C", new_x="LMARGIN", new_y="NEXT")
    pdf.cell(0, 8, f"Genere le {datetime.date.today().strftime('%d/%m/%Y')}", align="C", new_x="LMARGIN", new_y="NEXT")

    # Page 2: KPIs
    pdf.add_page()
    pdf.section_title("Indicateurs Cles de Performance")

    # Table header
    pdf.set_font("Helvetica", "B", 10)
    pdf.set_fill_color(192, 57, 43)
    pdf.set_text_color(255, 255, 255)
    pdf.cell(80, 8, "Indicateur", border=1, fill=True)
    pdf.cell(40, 8, "Valeur", border=1, fill=True)
    pdf.cell(40, 8, "Delta", border=1, fill=True, new_x="LMARGIN", new_y="NEXT")
    pdf.set_text_color(44, 62, 80)

    kpi_rows = [
        ("Top-of-Mind", f"{kpis['tom']}%", kpis['tom_delta']),
        ("Notoriete Totale", f"{kpis['notoriete']}%", kpis['notoriete_delta']),
        ("Penetration Parieurs", f"{kpis['penetration']}%", kpis['penetration_delta']),
        ("Consideration", f"{kpis['consideration']}%", kpis['consideration_delta']),
        ("Preference", f"{kpis['preference']}%", kpis['preference_delta']),
        ("Wallet Share", f"{kpis['wallet']}%", kpis['wallet_delta']),
        ("NPS Score", f"{kpis['nps']['nps']} pts", kpis['nps_delta']),
        ("Satisfaction", f"{kpis['satisfaction']}/5", kpis['satisfaction_delta']),
        ("Rappel Campagne", f"{kpis['rappel']}%", kpis['rappel_delta']),
    ]

    for i, (label, value, delta) in enumerate(kpi_rows):
        if i % 2 == 0:
            pdf.set_fill_color(245, 246, 248)
            fill = True
        else:
            fill = False

        pdf.set_font("Helvetica", "", 10)
        pdf.set_text_color(44, 62, 80)
        pdf.cell(80, 7, _safe(label), border=1, fill=fill)
        pdf.set_font("Helvetica", "B", 10)
        pdf.cell(40, 7, _safe(str(value)), border=1, fill=fill)

        pdf.set_font("Helvetica", "", 9)
        d_str = _safe(str(delta)) if delta else ""
        if d_str.startswith("+"):
            pdf.set_text_color(29, 131, 72)
        elif d_str.startswith("-"):
            pdf.set_text_color(192, 57, 43)
        else:
            pdf.set_text_color(212, 118, 10)
        pdf.cell(40, 7, d_str, border=1, fill=fill, new_x="LMARGIN", new_y="NEXT")

    # Page 3: Brand awareness
    pdf.add_page()
    pdf.section_title("Notoriete par Marque")

    pdf.set_font("Helvetica", "B", 10)
    pdf.set_fill_color(108, 52, 131)
    pdf.set_text_color(255, 255, 255)
    pdf.cell(60, 8, "Marque", border=1, fill=True)
    pdf.cell(50, 8, "Notoriete Totale", border=1, fill=True)
    pdf.cell(50, 8, "Top-of-Mind", border=1, fill=True, new_x="LMARGIN", new_y="NEXT")

    sorted_brands = sorted(kpis["notoriete_brands"].items(), key=lambda x: x[1], reverse=True)
    for i, (brand, not_val) in enumerate(sorted_brands):
        tom_val = kpis["tom_brands"].get(brand, 0)
        if brand == "Betclic":
            pdf.set_font("Helvetica", "B", 10)
            pdf.set_text_color(192, 57, 43)
        else:
            pdf.set_font("Helvetica", "", 10)
            pdf.set_text_color(44, 62, 80)

        fill = i % 2 == 0
        if fill:
            pdf.set_fill_color(245, 246, 248)

        pdf.cell(60, 7, _safe(brand), border=1, fill=fill)
        pdf.cell(50, 7, f"{not_val}%", border=1, fill=fill)
        pdf.cell(50, 7, f"{tom_val}%", border=1, fill=fill, new_x="LMARGIN", new_y="NEXT")

    # Page 4: Image
    pdf.add_page()
    pdf.section_title("Attributs Fonctionnels - Scores Moyens (/5)")

    pdf.set_font("Helvetica", "B", 10)
    pdf.set_fill_color(192, 57, 43)
    pdf.set_text_color(255, 255, 255)
    pdf.cell(90, 8, "Attribut", border=1, fill=True)
    pdf.cell(50, 8, "Score moyen", border=1, fill=True, new_x="LMARGIN", new_y="NEXT")

    sorted_attrs = sorted(kpis["image"].items(), key=lambda x: x[1], reverse=True)
    for i, (attr, score) in enumerate(sorted_attrs):
        pdf.set_font("Helvetica", "", 10)
        pdf.set_text_color(44, 62, 80)
        fill = i % 2 == 0
        if fill:
            pdf.set_fill_color(245, 246, 248)
        pdf.cell(90, 7, _safe(attr), border=1, fill=fill)
        pdf.cell(50, 7, f"{score}/5", border=1, fill=fill, new_x="LMARGIN", new_y="NEXT")

    buf = io.BytesIO()
    pdf.output(buf)
    buf.seek(0)
    return buf.getvalue()


# ═══════════════════════════════════════════
# EXCEL EXPORT
# ═══════════════════════════════════════════

def generate_excel(df, vagues_list) -> bytes:
    """Generate a comprehensive Excel report with multiple sheets."""
    kpis = _compute_all_kpis(df, vagues_list)

    buf = io.BytesIO()
    wb = xlsxwriter.Workbook(buf, {"in_memory": True})

    # Formats
    header_fmt = wb.add_format({
        "bold": True, "bg_color": "#C0392B", "font_color": "#FFFFFF",
        "border": 1, "font_size": 11, "align": "center", "valign": "vcenter",
    })
    header_purple_fmt = wb.add_format({
        "bold": True, "bg_color": "#6C3483", "font_color": "#FFFFFF",
        "border": 1, "font_size": 11, "align": "center", "valign": "vcenter",
    })
    title_fmt = wb.add_format({
        "bold": True, "font_size": 16, "font_color": "#C0392B",
    })
    subtitle_fmt = wb.add_format({
        "font_size": 11, "font_color": "#4A5568", "italic": True,
    })
    cell_fmt = wb.add_format({
        "border": 1, "font_size": 10, "align": "center", "valign": "vcenter",
    })
    cell_left_fmt = wb.add_format({
        "border": 1, "font_size": 10, "valign": "vcenter",
    })
    bold_cell_fmt = wb.add_format({
        "border": 1, "font_size": 10, "bold": True, "font_color": "#C0392B",
        "align": "center", "valign": "vcenter",
    })
    alt_row_fmt = wb.add_format({
        "border": 1, "font_size": 10, "bg_color": "#F5F6F8",
        "align": "center", "valign": "vcenter",
    })
    alt_row_left_fmt = wb.add_format({
        "border": 1, "font_size": 10, "bg_color": "#F5F6F8",
        "valign": "vcenter",
    })
    pct_fmt = wb.add_format({
        "border": 1, "font_size": 10, "num_format": "0.0%",
        "align": "center", "valign": "vcenter",
    })
    delta_up_fmt = wb.add_format({
        "border": 1, "font_size": 10, "font_color": "#1D8348",
        "align": "center", "valign": "vcenter",
    })
    delta_down_fmt = wb.add_format({
        "border": 1, "font_size": 10, "font_color": "#C0392B",
        "align": "center", "valign": "vcenter",
    })

    # ── Sheet 1: Scorecard ──
    ws1 = wb.add_worksheet("Scorecard KPIs")
    ws1.set_column("A:A", 28)
    ws1.set_column("B:B", 15)
    ws1.set_column("C:C", 18)

    ws1.write(0, 0, "Betclic Brand Pulse · Scorecard", title_fmt)
    ws1.write(1, 0, f"{kpis['latest_vague']} | Base : {kpis['n_total']} repondants", subtitle_fmt)

    row = 3
    ws1.write(row, 0, "Indicateur", header_fmt)
    ws1.write(row, 1, "Valeur", header_fmt)
    ws1.write(row, 2, "Delta vs precedent", header_fmt)

    kpi_rows = [
        ("Top-of-Mind", f"{kpis['tom']}%", kpis['tom_delta']),
        ("Notoriete Totale", f"{kpis['notoriete']}%", kpis['notoriete_delta']),
        ("Penetration Parieurs", f"{kpis['penetration']}%", kpis['penetration_delta']),
        ("Consideration", f"{kpis['consideration']}%", kpis['consideration_delta']),
        ("Preference", f"{kpis['preference']}%", kpis['preference_delta']),
        ("Wallet Share", f"{kpis['wallet']}%", kpis['wallet_delta']),
        ("NPS Score", f"{kpis['nps']['nps']} pts", kpis['nps_delta']),
        ("Satisfaction", f"{kpis['satisfaction']}/5", kpis['satisfaction_delta']),
        ("Rappel Campagne", f"{kpis['rappel']}%", kpis['rappel_delta']),
    ]

    for i, (label, value, delta) in enumerate(kpi_rows):
        r = row + 1 + i
        fmt_l = alt_row_left_fmt if i % 2 == 0 else cell_left_fmt
        fmt_c = alt_row_fmt if i % 2 == 0 else cell_fmt
        ws1.write(r, 0, label, fmt_l)
        ws1.write(r, 1, value, fmt_c)
        d_str = str(delta) if delta else ""
        d_fmt = delta_up_fmt if d_str.startswith("+") else (delta_down_fmt if d_str.startswith("-") else fmt_c)
        ws1.write(r, 2, d_str, d_fmt)

    # ── Sheet 2: Evolution par vague ──
    ws2 = wb.add_worksheet("Evolution Vagues")
    ws2.set_column("A:A", 28)
    ws2.set_column("B:D", 15)

    ws2.write(0, 0, "Evolution inter-vagues", title_fmt)
    row = 2
    ws2.write(row, 0, "Indicateur", header_fmt)
    _all_vagues = sorted([str(v) for v in df["Vague"].dropna().unique()],
                          key=lambda x: int(x.replace("Vague ", "")) if x.replace("Vague ", "").isdigit() else 999)
    for i, v in enumerate(_all_vagues):
        ws2.write(row, 1 + i, VAGUE_SHORT.get(v, v), header_fmt)

    evol_kpis = [
        ("Top-of-Mind", kpis["tom_vagues"]),
        ("Notoriete Totale", kpis["notoriete_vagues"]),
        ("Penetration", kpis["penetration_vagues"]),
        ("Consideration", kpis["consideration_vagues"]),
        ("Preference", kpis["preference_vagues"]),
        ("Wallet Share", kpis["wallet_vagues"]),
        ("Satisfaction", kpis["satisfaction_vagues"]),
        ("NPS", kpis["nps_vagues"]),
        ("Rappel Campagne", kpis["rappel_vagues"]),
    ]

    for i, (label, vague_data) in enumerate(evol_kpis):
        r = row + 1 + i
        fmt_l = alt_row_left_fmt if i % 2 == 0 else cell_left_fmt
        fmt_c = alt_row_fmt if i % 2 == 0 else cell_fmt
        ws2.write(r, 0, label, fmt_l)
        for j, v in enumerate(_all_vagues):
            val = vague_data.get(v)
            ws2.write(r, 1 + j, val if val is not None else "", fmt_c)

    # ── Sheet 3: Notoriete marques ──
    ws3 = wb.add_worksheet("Notoriete Marques")
    ws3.set_column("A:A", 20)
    ws3.set_column("B:C", 18)

    ws3.write(0, 0, "Notoriete par marque", title_fmt)
    row = 2
    ws3.write(row, 0, "Marque", header_purple_fmt)
    ws3.write(row, 1, "Notoriete Totale", header_purple_fmt)
    ws3.write(row, 2, "Top-of-Mind", header_purple_fmt)

    sorted_brands = sorted(kpis["notoriete_brands"].items(), key=lambda x: x[1], reverse=True)
    for i, (brand, not_val) in enumerate(sorted_brands):
        r = row + 1 + i
        tom_val = kpis["tom_brands"].get(brand, 0)
        if brand == "Betclic":
            ws3.write(r, 0, brand, bold_cell_fmt)
            ws3.write(r, 1, f"{not_val}%", bold_cell_fmt)
            ws3.write(r, 2, f"{tom_val}%", bold_cell_fmt)
        else:
            fmt_l = alt_row_left_fmt if i % 2 == 0 else cell_left_fmt
            fmt_c = alt_row_fmt if i % 2 == 0 else cell_fmt
            ws3.write(r, 0, brand, fmt_l)
            ws3.write(r, 1, f"{not_val}%", fmt_c)
            ws3.write(r, 2, f"{tom_val}%", fmt_c)

    # ── Sheet 4: Image de marque ──
    ws4 = wb.add_worksheet("Attributs Fonctionnels")
    ws4.set_column("A:A", 28)
    ws4.set_column("B:B", 15)

    ws4.write(0, 0, "Attributs Fonctionnels · Scores (/5)", title_fmt)
    row = 2
    ws4.write(row, 0, "Attribut", header_fmt)
    ws4.write(row, 1, "Score moyen", header_fmt)

    sorted_attrs = sorted(kpis["image"].items(), key=lambda x: x[1], reverse=True)
    for i, (attr, score) in enumerate(sorted_attrs):
        r = row + 1 + i
        fmt_l = alt_row_left_fmt if i % 2 == 0 else cell_left_fmt
        fmt_c = alt_row_fmt if i % 2 == 0 else cell_fmt
        ws4.write(r, 0, attr, fmt_l)
        ws4.write(r, 1, score, fmt_c)

    # ── Sheet 5: Funnel ──
    ws5 = wb.add_worksheet("Funnel")
    ws5.set_column("A:A", 22)
    ws5.set_column("B:B", 12)

    ws5.write(0, 0, "Funnel Betclic", title_fmt)
    row = 2
    ws5.write(row, 0, "Etape", header_fmt)
    ws5.write(row, 1, "Valeur", header_fmt)

    for i, (step, val) in enumerate(kpis["funnel"].items()):
        r = row + 1 + i
        fmt_l = alt_row_left_fmt if i % 2 == 0 else cell_left_fmt
        fmt_c = alt_row_fmt if i % 2 == 0 else cell_fmt
        ws5.write(r, 0, step, fmt_l)
        ws5.write(r, 1, f"{val}%", fmt_c)

    wb.close()
    buf.seek(0)
    return buf.getvalue()
